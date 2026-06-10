#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 生成工具 - 支持中文
基于 python-pptx 库
"""

import os
from typing import List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    from .fonts import font_manager, TEST_STRINGS
except ImportError:
    from fonts import font_manager, TEST_STRINGS


class PPTGenerator:
    """PPT 生成器"""
    
    def __init__(self):
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")
        self.font_name = font_manager.get_font_name()
    
    def create_presentation(self, output_path: str, slides: List[dict]) -> str:
        """
        创建演示文稿
        
        参数:
            output_path: 输出路径
            slides: 幻灯片列表，每个幻灯片是字典:
                {
                    'title': '标题',
                    'content': '内容文本或列表',
                    'layout': 'title' | 'title_and_content' | 'blank'
                }
        """
        prs = Presentation()
        
        for slide_data in slides:
            layout_type = slide_data.get('layout', 'title_and_content')
            
            if layout_type == 'title':
                slide_layout = prs.slide_layouts[0]  # 标题幻灯片
            elif layout_type == 'blank':
                slide_layout = prs.slide_layouts[6]  # 空白幻灯片
            else:
                slide_layout = prs.slide_layouts[1]  # 标题和内容
            
            slide = prs.slides.add_slide(slide_layout)
            
            # 标题
            title = slide_data.get('title', '')
            if title and slide.shapes.title:
                title_box = slide.shapes.title
                title_box.text = title
                # 设置字体
                for paragraph in title_box.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_name
                        run.font.size = Pt(32)
                        run.font.bold = True
            
            # 内容
            content = slide_data.get('content', '')
            if content:
                # 查找内容占位符
                if len(slide.placeholders) > 1:
                    content_box = slide.placeholders[1]
                else:
                    # 添加文本框
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(5)
                    content_box = slide.shapes.add_textbox(left, top, width, height)
                
                tf = content_box.text_frame
                tf.text = content if isinstance(content, str) else '\n'.join(content)
                
                # 设置字体
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_name
                        run.font.size = Pt(18)
        
        prs.save(output_path)
        return output_path
    
    def create_title_slide(self, output_path: str, main_title: str, 
                          subtitle: str = "") -> str:
        """创建标题幻灯片"""
        slides = [{
            'title': main_title,
            'content': subtitle,
            'layout': 'title'
        }]
        return self.create_presentation(output_path, slides)
    
    def create_content_slide(self, output_path: str, title: str, 
                            content: List[str]) -> str:
        """创建内容幻灯片"""
        slides = [{
            'title': title,
            'content': content,
            'layout': 'title_and_content'
        }]
        return self.create_presentation(output_path, slides)


def create_simple_ppt(output_path: str, title: str, 
                      content_slides: List[Tuple[str, str]]) -> str:
    """
    快速创建简单PPT
    
    参数:
        output_path: 输出路径
        title: 标题幻灯片标题
        content_slides: [(slide_title, slide_content), ...]
    
    例子:
        create_simple_ppt('test.pptx', '考研复习', [
            ('第一章', '细胞的基本功能'),
            ('第二章', '神经传导'),
        ])
    """
    generator = PPTGenerator()
    
    slides = [{'title': title, 'content': '', 'layout': 'title'}]
    
    for slide_title, slide_content in content_slides:
        slides.append({
            'title': slide_title,
            'content': slide_content,
            'layout': 'title_and_content'
        })
    
    return generator.create_presentation(output_path, slides)


def create_test_ppt(output_path: str) -> str:
    """创建测试PPT，验证中文显示"""
    generator = PPTGenerator()
    
    slides = [
        {
            'title': '中文PPT测试',
            'content': '验证各种字符的显示效果',
            'layout': 'title'
        },
        {
            'title': '中文内容测试',
            'content': TEST_STRINGS['chinese'],
            'layout': 'title_and_content'
        },
        {
            'title': '标点符号测试',
            'content': TEST_STRINGS['punctuation'],
            'layout': 'title_and_content'
        },
        {
            'title': '数学符号测试',
            'content': TEST_STRINGS['math'],
            'layout': 'title_and_content'
        },
        {
            'title': '特殊符号测试',
            'content': TEST_STRINGS['symbols'],
            'layout': 'title_and_content'
        },
        {
            'title': '箭头符号测试',
            'content': TEST_STRINGS['arrows'],
            'layout': 'title_and_content'
        },
        {
            'title': '货币符号测试',
            'content': TEST_STRINGS['currency'],
            'layout': 'title_and_content'
        },
        {
            'title': '混合内容测试',
            'content': TEST_STRINGS['mixed'],
            'layout': 'title_and_content'
        },
    ]
    
    return generator.create_presentation(output_path, slides)


if __name__ == "__main__":
    print("测试PPT生成...")
    try:
        path = create_test_ppt('/tmp/test_chinese.pptx')
        print(f"✓ 测试PPT已生成: {path}")
    except Exception as e:
        print(f"✗ 测试失败: {e}")
