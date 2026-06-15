#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - 支持中文内容
基于 python-pptx 库
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .fonts import get_font_name
except ImportError:
    from fonts import get_font_name


class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self):
        self.font_name = get_font_name()
    
    def create_presentation(self, output_path: str, slides: list) -> str:
        """
        创建演示文稿
        
        参数:
            output_path: 输出路径
            slides: 幻灯片列表，每个是字典
                {
                    'title': '标题',
                    'content': '内容文本或列表',
                    'layout': 'title' | 'title_and_content' | 'blank'
                }
        """
        prs = Presentation()
        
        for slide_data in slides:
            layout_type = slide_data.get('layout', 'title_and_content')
            
            if layout_type == 'title':
                slide_layout = prs.slide_layouts[0]
            elif layout_type == 'blank':
                slide_layout = prs.slide_layouts[6]
            else:
                slide_layout = prs.slide_layouts[1]
            
            slide = prs.slides.add_slide(slide_layout)
            
            # 标题
            title = slide_data.get('title', '')
            if title and slide.shapes.title:
                title_box = slide.shapes.title
                title_box.text = title
                for paragraph in title_box.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_name
                        run.font.size = Pt(32)
                        run.font.bold = True
            
            # 内容
            content = slide_data.get('content', '')
            if content:
                if len(slide.placeholders) > 1:
                    content_box = slide.placeholders[1]
                else:
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(5)
                    content_box = slide.shapes.add_textbox(left, top, width, height)
                
                tf = content_box.text_frame
                tf.text = content if isinstance(content, str) else '\n'.join(content)
                
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font_name
                        run.font.size = Pt(18)
        
        prs.save(output_path)
        return output_path


def create_simple_ppt(output_path: str, title: str, content_slides: list) -> str:
    """
    快速创建简单PPT
    
    参数:
        output_path: 输出路径
        title: 标题幻灯片标题
        content_slides: [(slide_title, slide_content), ...]
    """
    generator = PPTGenerator()
    
    slides = [{'title': title, 'content': '', 'layout': 'title'}]
    
    for slide_title, slide_content in content_slides:
        slides.append({
            'title': slide_title,
            'content': slide_content,
            'layout': 'title_and_content'
        })
    
    return generator.create_presentation(output_path, slides)
