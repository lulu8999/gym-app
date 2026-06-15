#!/usr/bin/env python3
"""
Markdown 转 Office 套件（Word/PPT/Excel）完整解决方案
支持中文、表格、图片等复杂内容
"""

import subprocess
import tempfile
import os
import re
from typing import Optional, List, Dict, Any
from pathlib import Path

# Word 转换
def md_to_docx(
    content: str,
    output_path: str,
    title: Optional[str] = None,
    toc: bool = True,
    extra_args: Optional[List[str]] = None
) -> str:
    """
    Markdown 转 Word 文档

    Args:
        content: Markdown 内容
        output_path: 输出路径
        title: 文档标题
        toc: 是否生成目录
        extra_args: 额外参数

    Returns:
        输出文件路径
    """
    # 检查 pandoc
    try:
        subprocess.run(['pandoc', '--version'], capture_output=True, check=True)
    except (subprocess.CalledProcessError, FileNotFoundError):
        raise RuntimeError("未找到 pandoc，请先安装: dnf install pandoc")

    # 添加标题
    if title:
        content = f"# {title}\n\n{content}"

    # 创建临时 markdown 文件
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as f:
        f.write(content)
        md_path = f.name

    try:
        # 构建命令
        cmd = ['pandoc', md_path, '-o', output_path, '--from', 'markdown', '--to', 'docx']

        if toc:
            cmd.append('--toc')

        if extra_args:
            cmd.extend(extra_args)

        # 执行转换
        result = subprocess.run(cmd, capture_output=True, text=True)

        if result.returncode != 0:
            raise RuntimeError(f"Pandoc 转换失败: {result.stderr}")

        return output_path

    finally:
        os.unlink(md_path)


# PPT 转换
def md_to_pptx(
    content: str,
    output_path: str,
    title: Optional[str] = None,
    subtitle: Optional[str] = None
) -> str:
    """
    Markdown 转 PPT 演示文稿

    格式约定：
    - # 标题 → 新幻灯片标题
    - ## 副标题 → 幻灯片内容标题
    - - 列表项 → 幻灯片正文内容
    - > 引用 → 重点提示

    Args:
        content: Markdown 内容
        output_path: 输出路径
        title: 总标题（第一张幻灯片）
        subtitle: 副标题

    Returns:
        输出文件路径
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 解析 Markdown 内容
    slides_data = _parse_md_for_ppt(content, title, subtitle)

    for slide_data in slides_data:
        if slide_data['type'] == 'title':
            # 标题幻灯片
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get('title', '')
            if slide.placeholders[1]:
                slide.placeholders[1].text = slide_data.get('subtitle', '')

        elif slide_data['type'] == 'content':
            # 内容幻灯片
            slide_layout = prs.slide_layouts[1]  # 标题和内容布局
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get('title', '')

            # 添加内容
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True

            for item in slide_data.get('content', []):
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)

    prs.save(output_path)
    return output_path


def _parse_md_for_ppt(content: str, main_title: Optional[str], main_subtitle: Optional[str]) -> List[Dict]:
    """
    解析 Markdown 内容为 PPT 结构
    """
    slides = []
    current_slide = None

    # 如果有主标题，创建标题幻灯片
    if main_title:
        slides.append({
            'type': 'title',
            'title': main_title,
            'subtitle': main_subtitle or ''
        })

    lines = content.strip().split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()

        # 一级标题 = 新幻灯片
        if line.startswith('# ') and not line.startswith('##'):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'content',
                'title': line[2:].strip(),
                'content': []
            }

        # 二级标题 = 幻灯片内小标题
        elif line.startswith('## '):
            if current_slide:
                current_slide['content'].append(line[3:].strip())

        # 列表项
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide['content'].append('• ' + line[2:].strip())

        # 引用
        elif line.startswith('> '):
            if current_slide:
                current_slide['content'].append('【重点】' + line[2:].strip())

        # 空行，可能表示分隔
        elif not line and current_slide and current_slide['content']:
            pass

        i += 1

    # 添加最后一张幻灯片
    if current_slide:
        slides.append(current_slide)

    return slides


# Excel 转换
def md_to_excel(
    content: str,
    output_path: str,
    sheet_name: str = 'Sheet1'
) -> str:
    """
    Markdown 表格 转 Excel

    支持多个表格，每个表格会生成一个工作表

    Args:
        content: Markdown 内容（包含表格）
        output_path: 输出路径
        sheet_name: 默认工作表名称

    Returns:
        输出文件路径
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

    wb = Workbook()

    # 解析所有表格
    tables = _extract_tables_from_md(content)

    if not tables:
        # 没有表格，创建空工作簿
        ws = wb.active
        ws.title = sheet_name
    else:
        for idx, table in enumerate(tables):
            if idx == 0:
                ws = wb.active
                ws.title = f"表格{idx+1}" if len(tables) > 1 else sheet_name
            else:
                ws = wb.create_sheet(title=f"表格{idx+1}")

            # 填充数据
            for row_idx, row_data in enumerate(table, 1):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)

                    # 第一行加粗（表头）
                    if row_idx == 1:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
                        cell.font = Font(bold=True, color="FFFFFF")

                    # 居中对齐
                    cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)

            # 调整列宽
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column].width = adjusted_width

    wb.save(output_path)
    return output_path


def _extract_tables_from_md(content: str) -> List[List[List[str]]]:
    """
    从 Markdown 中提取所有表格

    Returns:
        列表，每个表格是一个二维列表
    """
    tables = []
    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # 检测表格开始（包含 | 的行）
        if '|' in line and not line.startswith('>'):
            table_lines = []
            while i < len(lines) and '|' in lines[i]:
                table_lines.append(lines[i])
                i += 1

            # 解析表格
            if len(table_lines) >= 2:  # 至少需要标题行和分隔行
                table_data = []
                for line in table_lines:
                    # 跳过分隔行（包含 ---）
                    if '---' in line.replace('|', '').replace('-', '').strip():
                        continue

                    # 解析单元格
                    cells = [cell.strip() for cell in line.split('|')]
                    # 过滤空单元格（开头和结尾的）
                    cells = [c for c in cells if c]
                    if cells:
                        table_data.append(cells)

                if table_data:
                    tables.append(table_data)
        else:
            i += 1

    return tables


# 统一入口函数
def convert_md(
    content: str,
    output_path: str,
    format: str = 'docx',
    **kwargs
) -> str:
    """
    统一转换入口

    Args:
        content: Markdown 内容
        output_path: 输出路径
        format: 输出格式 ('docx', 'pptx', 'xlsx')
        **kwargs: 额外参数

    Returns:
        输出文件路径
    """
    format = format.lower()

    if format == 'docx':
        return md_to_docx(content, output_path, **kwargs)
    elif format == 'pptx':
        return md_to_pptx(content, output_path, **kwargs)
    elif format == 'xlsx':
        return md_to_excel(content, output_path, **kwargs)
    else:
        raise ValueError(f"不支持的格式: {format}")


def validate_docx(docx_path: str) -> dict:
    """
    验证 Word 文档质量

    Args:
        docx_path: Word 文件路径

    Returns:
        验证结果字典
    """
    from docx import Document

    doc = Document(docx_path)

    # 统计信息
    total_paras = len(doc.paragraphs)
    total_tables = len(doc.tables)
    empty_lines = sum(1 for p in doc.paragraphs if not p.text.strip())
    short_paras = sum(1 for p in doc.paragraphs if len(p.text.strip()) < 20)

    # 检测问题
    issues = []
    if total_paras > 300:
        issues.append(f"段落数过多({total_paras})，可能是表格被拆散")
    if empty_lines > 20:
        issues.append(f"空行过多({empty_lines})")
    if short_paras > 100:
        issues.append(f"短段落过多({short_paras})，可能是表格内容被拆散")

    return {
        'valid': len(issues) == 0,
        'paragraphs': total_paras,
        'tables': total_tables,
        'empty_lines': empty_lines,
        'short_paragraphs': short_paras,
        'issues': issues
    }


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 4:
        print("用法: python md_to_office.py <format> <input.md> <output>")
        print("")
        print("格式:")
        print("  docx  - Word 文档")
        print("  pptx  - PPT 演示文稿")
        print("  xlsx  - Excel 表格")
        sys.exit(1)

    fmt = sys.argv[1].lower()
    input_file = sys.argv[2]
    output_file = sys.argv[3]

    # 读取输入文件
    with open(input_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # 执行转换
    result = convert_md(content, output_file, format=fmt)
    print(f"✓ 转换成功: {result}")
